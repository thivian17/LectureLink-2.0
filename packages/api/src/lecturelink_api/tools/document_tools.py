"""Document text extraction tool for the LectureLink syllabus processor.

Routes extraction by file type:
- PDF  → Gemini 2.5 Flash multimodal (via google-genai SDK)
- DOCX → python-docx local extraction
- PPTX → python-pptx local extraction
"""

from __future__ import annotations

import asyncio
import io
from typing import Any

from google import genai
from google.genai import types
from loguru import logger

from lecturelink_api.config import get_settings

# ---------------------------------------------------------------------------
# PDF extraction via Gemini 2.5 Flash
# ---------------------------------------------------------------------------

async def _extract_pdf_with_gemini(file_bytes: bytes) -> dict[str, Any]:
    """Send raw PDF bytes to Gemini 2.5 Flash and ask for structured text extraction."""
    try:
        settings = get_settings()
        if settings.GOOGLE_API_KEY:
            client = genai.Client(api_key=settings.GOOGLE_API_KEY)
        else:
            client = genai.Client(
                vertexai=True,
                project=settings.GOOGLE_CLOUD_PROJECT,
                location="us-central1",
            )

        pdf_part = types.Part.from_bytes(data=file_bytes, mime_type="application/pdf")

        prompt = (
            "Extract ALL text from this PDF document. "
            "Preserve the original structure including headings, sub-headings, "
            "bullet/numbered lists, and tables. "
            "For tables, reproduce them as markdown tables. "
            "Return only the extracted text with no additional commentary."
        )

        response = await asyncio.to_thread(
            client.models.generate_content,
            model="gemini-2.5-flash",
            contents=[pdf_part, prompt],
        )

        text = response.text or ""
        return {
            "text": text,
            "method": "gemini-pdf",
            "page_count": None,
        }
    except Exception as exc:
        logger.error("Gemini PDF extraction failed: {}", exc)
        return {
            "text": "",
            "method": "gemini-pdf",
            "page_count": None,
            "error": f"Gemini extraction failed: {exc}",
        }


# ---------------------------------------------------------------------------
# DOCX extraction via python-docx
# ---------------------------------------------------------------------------

_HEADING_MARKERS: dict[str, str] = {
    "Heading 1": "# ",
    "Heading 2": "## ",
    "Heading 3": "### ",
    "Heading 4": "#### ",
    "Title": "# ",
    "Subtitle": "## ",
}

_SKIP_STYLES: set[str] = {"toc 1", "toc 2", "toc 3", "toc heading"}


def _get_heading_prefix(style_name: str) -> str | None:
    """Get markdown heading prefix for a paragraph style.

    Returns None if the style should be skipped (e.g., TOC entries).
    Returns "" for normal body text.
    Returns "# ", "## ", etc. for heading styles.

    Handles custom heading styles by checking for 'Heading' keyword
    and extracting the level number.
    """
    if not style_name:
        return ""

    if style_name.lower() in _SKIP_STYLES:
        return None

    if style_name in _HEADING_MARKERS:
        return _HEADING_MARKERS[style_name]

    style_lower = style_name.lower()
    if "heading" in style_lower:
        import re

        match = re.search(r"heading\s*(\d)", style_lower)
        if match:
            level = min(int(match.group(1)), 4)
            return "#" * level + " "
        return "## "

    if "title" in style_lower:
        return "# "
    if "subtitle" in style_lower:
        return "## "

    return ""


def _extract_textbox_content(element) -> list[str]:
    """Extract text from Word text boxes (w:txbxContent) nested within an element.

    Text boxes in DOCX appear inside structures like:
      mc:AlternateContent > mc:Choice > w:drawing > ... > wps:txbx > w:txbxContent > w:p

    This function recursively searches for all w:txbxContent elements and
    extracts paragraph text from each.
    """
    from docx.oxml.ns import qn

    texts = []
    for txbx_content in element.iter(qn("w:txbxContent")):
        for p_elem in txbx_content.iterchildren(qn("w:p")):
            runs_text = []
            for r_elem in p_elem.iter(qn("w:t")):
                if r_elem.text:
                    runs_text.append(r_elem.text)
            text = "".join(runs_text).strip()
            if text:
                texts.append(text)
    return texts


def _table_to_markdown(table) -> str:
    """Convert a python-docx Table to a markdown-formatted table string.

    Handles merged cells by deduplicating adjacent cells with identical text
    (Word reports the same text for each grid column a merged cell spans).
    """
    rows: list[list[str]] = []
    for row in table.rows:
        cells_text = [cell.text.strip() for cell in row.cells]
        # Deduplicate adjacent identical cells (merged cells)
        deduped: list[str] = []
        prev = None
        for text in cells_text:
            if text != prev:
                deduped.append(text)
            prev = text
        rows.append(deduped)

    if not rows:
        return ""

    # Normalize column count (merged cells may produce different lengths per row)
    max_cols = max(len(r) for r in rows)
    for r in rows:
        while len(r) < max_cols:
            r.append("")

    lines: list[str] = []
    lines.append("| " + " | ".join(rows[0]) + " |")
    lines.append("| " + " | ".join("---" for _ in rows[0]) + " |")
    for row in rows[1:]:
        lines.append("| " + " | ".join(row) + " |")

    return "\n".join(lines)


async def _extract_docx(file_bytes: bytes) -> dict[str, Any]:
    """Extract text from a DOCX file preserving headings and tables."""
    from docx import Document

    def _do_extract() -> dict[str, Any]:
        doc = Document(io.BytesIO(file_bytes))
        parts: list[str] = []
        seen_textbox_texts: set[str] = set()  # deduplicate mc:Choice + mc:Fallback

        from docx.oxml.ns import qn

        # --- Extract header/footer content (often contains course code, dept) ---
        header_footer_texts: list[str] = []
        seen_hf: set[str] = set()

        for section in doc.sections:
            try:
                for para in section.header.paragraphs:
                    text = para.text.strip()
                    if text and text not in seen_hf:
                        seen_hf.add(text)
                        header_footer_texts.append(text)
            except Exception:
                pass
            try:
                for para in section.footer.paragraphs:
                    text = para.text.strip()
                    if text and text not in seen_hf:
                        seen_hf.add(text)
                        header_footer_texts.append(text)
            except Exception:
                pass

        if header_footer_texts:
            parts.append("--- Document Header/Footer ---")
            parts.extend(header_footer_texts)
            parts.append("--- End Header/Footer ---")

        # --- Iterate over body elements preserving document order ---
        for element in doc.element.body:
            tag = element.tag

            if tag == qn("w:p"):
                from docx.text.paragraph import Paragraph

                para = Paragraph(element, doc)
                style_name = para.style.name if para.style else "Normal"

                prefix = _get_heading_prefix(style_name)
                if prefix is None:
                    continue  # Skip TOC entries

                # Bullet prefix for list paragraphs
                if style_name == "List Paragraph":
                    prefix = "- "

                text = para.text.strip()
                if text:
                    parts.append(f"{prefix}{text}")

            elif tag == qn("w:tbl"):
                from docx.table import Table

                tbl = Table(element, doc)
                md = _table_to_markdown(tbl)
                if md:
                    parts.append(md)

            else:
                # Handle text boxes inside mc:AlternateContent, w:drawing, etc.
                textbox_texts = _extract_textbox_content(element)
                for txt in textbox_texts:
                    if txt not in seen_textbox_texts:
                        seen_textbox_texts.add(txt)
                        parts.append(txt)

        return {
            "text": "\n\n".join(parts),
            "method": "docx",
            "page_count": None,
        }

    return await asyncio.to_thread(_do_extract)


# ---------------------------------------------------------------------------
# PPTX extraction via python-pptx
# ---------------------------------------------------------------------------

def _shapes_text(shapes) -> list[str]:
    """Recursively extract text from shapes, including group shapes."""
    texts: list[str] = []
    for shape in shapes:
        if shape.shape_type == 6:  # MSO_SHAPE_TYPE.GROUP
            texts.extend(_shapes_text(shape.shapes))
        elif shape.has_table:
            rows: list[list[str]] = []
            for row in shape.table.rows:
                rows.append([cell.text.strip() for cell in row.cells])
            if rows:
                lines = []
                lines.append("| " + " | ".join(rows[0]) + " |")
                lines.append("| " + " | ".join("---" for _ in rows[0]) + " |")
                for row in rows[1:]:
                    lines.append("| " + " | ".join(row) + " |")
                texts.append("\n".join(lines))
        elif shape.has_text_frame:
            frame_text = "\n".join(
                para.text.strip() for para in shape.text_frame.paragraphs if para.text.strip()
            )
            if frame_text:
                texts.append(frame_text)
    return texts


async def _extract_pptx(file_bytes: bytes) -> dict[str, Any]:
    """Extract text from a PPTX file including titles, body text, tables, and notes."""
    from pptx import Presentation

    def _do_extract() -> dict[str, Any]:
        prs = Presentation(io.BytesIO(file_bytes))
        slide_parts: list[str] = []

        for idx, slide in enumerate(prs.slides, start=1):
            # Title
            title = ""
            if slide.shapes.title and slide.shapes.title.has_text_frame:
                title = slide.shapes.title.text_frame.text.strip()

            header = f"## Slide {idx}: {title}" if title else f"## Slide {idx}"

            # Body content from all shapes
            content_lines = _shapes_text(slide.shapes)
            content = "\n".join(content_lines)

            # Speaker notes
            notes = ""
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes = slide.notes_slide.notes_text_frame.text.strip()

            section = header
            if content:
                section += f"\n{content}"
            if notes:
                section += f"\n\n**Notes:** {notes}"

            slide_parts.append(section)

        return {
            "text": "\n\n".join(slide_parts),
            "method": "pptx",
            "page_count": len(prs.slides),
        }

    return await asyncio.to_thread(_do_extract)


# ---------------------------------------------------------------------------
# Public API
# ---------------------------------------------------------------------------

_MIME_ROUTER: dict[str, Any] = {
    "application/pdf": _extract_pdf_with_gemini,
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document": _extract_docx,
    "application/vnd.openxmlformats-officedocument.presentationml.presentation": _extract_pptx,
}


async def extract_document_text(
    file_bytes: bytes,
    file_name: str,
    mime_type: str,
) -> dict[str, Any]:
    """Extract text from an uploaded syllabus document.

    Args:
        file_bytes: Raw bytes of the uploaded file.
        file_name: Original filename (used for logging).
        mime_type: MIME type of the file.

    Returns:
        Dict with keys: text, method, page_count, and optionally error.
    """
    logger.info("Extracting text from '{}' ({})", file_name, mime_type)

    handler = _MIME_ROUTER.get(mime_type)
    if handler is None:
        return {
            "text": "",
            "method": "unsupported",
            "page_count": None,
            "error": f"Unsupported file type: {mime_type}",
        }

    return await handler(file_bytes)
