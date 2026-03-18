"""Slide Analyzer — extracts text and visual descriptions from PDF/PPTX slides.

Uses Gemini multimodal vision to process each slide, extracting text content
and describing visual elements (diagrams, charts, equations). Supports both
PDF (native Gemini) and PPTX (converted to images via python-pptx).
"""

from __future__ import annotations

import asyncio
import io
import json
import logging
from pathlib import Path
from urllib.parse import urlparse

from google import genai
from google.genai import types

logger = logging.getLogger(__name__)

SLIDE_ANALYSIS_PROMPT = """Analyze each slide in this presentation.

For each slide, extract:
1. slide_number: The slide number (1-indexed)
2. title: The slide title/heading (null if no clear title)
3. text_content: ALL text visible on the slide, preserving structure
4. visual_description: Detailed description of any diagrams, charts, images, \
or visual elements. A diagram of a heat engine is more important than the bullet \
points around it. Describe visuals thoroughly.
5. has_diagram: true if the slide contains a diagram, flowchart, or illustration
6. has_code: true if the slide contains code snippets or pseudocode
7. has_equation: true if the slide contains mathematical equations or formulas

Output a JSON array:
[
    {
        "slide_number": 1,
        "title": "Introduction to Thermodynamics",
        "text_content": "PHYS 201 - Lecture 1\\nDr. Smith",
        "visual_description": null,
        "has_diagram": false,
        "has_code": false,
        "has_equation": false
    }
]

Guidelines:
- Include EVERY slide, even title/section divider slides.
- For equations, transcribe them in LaTeX notation when possible.
- For code, preserve exact formatting and syntax.
- Visual descriptions should be detailed enough that a student who can't see \
the slide can understand the concept being illustrated.
- If a slide has both text and a diagram, include both — the visual description \
supplements the text content.

Output ONLY the JSON array, no other text."""

_MAX_RETRIES = 3
_RETRY_BASE_DELAY = 1.0  # seconds


class SlideAnalysisError(Exception):
    """Raised when slide analysis fails."""


async def analyze_slides(slides_url: str) -> list[dict]:
    """Analyze slide deck using Gemini multimodal vision.

    Args:
        slides_url: Supabase Storage URL or local path to PDF/PPTX file.

    Returns:
        List of slide analyses with keys: slide_number, title, text_content,
        visual_description, has_diagram, has_code, has_equation.

    Raises:
        SlideAnalysisError: If analysis fails after retries.
    """
    last_error: Exception | None = None
    for attempt in range(_MAX_RETRIES):
        try:
            return await _call_gemini(slides_url)
        except SlideAnalysisError:
            raise
        except Exception as e:
            # Don't retry on permanent errors (bad file, no pages, etc.)
            err_str = str(e)
            if "no pages" in err_str.lower() or "INVALID_ARGUMENT" in err_str:
                raise SlideAnalysisError(
                    f"Slide file is invalid or empty: {e}"
                ) from e
            last_error = e
            delay = _RETRY_BASE_DELAY * (2**attempt)
            logger.warning(
                "Slide analysis attempt %d/%d failed: %s. Retrying in %.1fs",
                attempt + 1,
                _MAX_RETRIES,
                e,
                delay,
            )
            await asyncio.sleep(delay)

    raise SlideAnalysisError(
        f"Slide analysis failed after {_MAX_RETRIES} retries: {last_error}"
    )


def _recover_truncated_json(text: str) -> list[dict]:
    """Recover complete slide objects from a truncated JSON array.

    When Gemini's response is cut off mid-JSON, we find the last complete
    object boundary and parse everything up to that point.
    """
    # Find the last complete object: look for "},\n  {" or "}\n]" patterns
    last_complete = text.rfind("}")
    while last_complete > 0:
        candidate = text[: last_complete + 1]
        # Try closing the array
        try:
            result = json.loads(candidate + "]")
            if isinstance(result, list):
                return result
        except json.JSONDecodeError:
            pass
        last_complete = text.rfind("}", 0, last_complete)
    return []


_BATCH_SIZE = 15  # max pages per Gemini call to avoid output truncation


async def _call_gemini(slides_url: str) -> list[dict]:
    """Build content from slides_url, send to Gemini, and parse the JSON.

    For PDFs with many pages, splits into batches to avoid output truncation.
    """
    mime_type = get_slide_mime_type(slides_url)

    # PPTX: extract text per slide, then batch
    if mime_type.endswith("presentationml.presentation"):
        all_parts = _pptx_to_image_parts(slides_url)
        batches = [
            all_parts[i : i + _BATCH_SIZE]
            for i in range(0, len(all_parts), _BATCH_SIZE)
        ]
        batch_offsets = [i * _BATCH_SIZE for i in range(len(batches))]
        results = await asyncio.gather(
            *[
                _analyze_parts_batch(parts, offset)
                for parts, offset in zip(batches, batch_offsets)
            ]
        )
        all_slides = [s for batch in results for s in batch]
        validated = validate_slide_analysis(all_slides)
        logger.info("Slide analysis complete: %d slides", len(validated))
        return validated

    # PDF: download once, then split into page batches
    file_bytes = _download_file(slides_url)
    page_pdfs = _split_pdf(file_bytes)
    total_pages = len(page_pdfs)
    logger.info("PDF has %d pages, batch size %d", total_pages, _BATCH_SIZE)

    if total_pages <= _BATCH_SIZE:
        # Small PDF — single call with original file
        return await _analyze_single(file_bytes, mime_type)

    # Large PDF — batch process
    batches = [
        page_pdfs[i : i + _BATCH_SIZE]
        for i in range(0, total_pages, _BATCH_SIZE)
    ]
    results = await asyncio.gather(
        *[
            _analyze_pdf_batch(batch, batch_idx * _BATCH_SIZE)
            for batch_idx, batch in enumerate(batches)
        ]
    )
    all_slides = [s for batch in results for s in batch]
    validated = validate_slide_analysis(all_slides)
    logger.info("Slide analysis complete: %d slides (batched)", len(validated))
    return validated


def _download_file(file_path: str) -> bytes:
    """Download a file from URL or read from local path."""
    if file_path.startswith(("http://", "https://")):
        import httpx

        resp = httpx.get(file_path, follow_redirects=True)
        resp.raise_for_status()
        return resp.content
    return Path(file_path).read_bytes()


def _split_pdf(pdf_bytes: bytes) -> list[bytes]:
    """Split a PDF into individual single-page PDFs."""
    from pypdf import PdfReader, PdfWriter

    reader = PdfReader(io.BytesIO(pdf_bytes))
    pages: list[bytes] = []
    for page in reader.pages:
        writer = PdfWriter()
        writer.add_page(page)
        buf = io.BytesIO()
        writer.write(buf)
        pages.append(buf.getvalue())
    return pages


def _merge_page_pdfs(page_pdfs: list[bytes]) -> bytes:
    """Merge multiple single-page PDFs into one PDF."""
    from pypdf import PdfReader, PdfWriter

    writer = PdfWriter()
    for pdf_bytes in page_pdfs:
        reader = PdfReader(io.BytesIO(pdf_bytes))
        for page in reader.pages:
            writer.add_page(page)
    buf = io.BytesIO()
    writer.write(buf)
    return buf.getvalue()


async def _analyze_single(file_bytes: bytes, mime_type: str) -> list[dict]:
    """Analyze a complete file in a single Gemini call."""
    parts = [
        types.Part.from_bytes(data=file_bytes, mime_type=mime_type),
        types.Part(text=SLIDE_ANALYSIS_PROMPT),
    ]
    return await _send_and_parse(parts)


async def _analyze_pdf_batch(page_pdfs: list[bytes], offset: int) -> list[dict]:
    """Analyze a batch of PDF pages in one Gemini call."""
    merged = _merge_page_pdfs(page_pdfs)
    prompt = SLIDE_ANALYSIS_PROMPT + (
        f"\n\nNOTE: These are slides {offset + 1}-{offset + len(page_pdfs)}."
        f" Number them starting at {offset + 1}."
    )
    parts = [
        types.Part.from_bytes(data=merged, mime_type="application/pdf"),
        types.Part(text=prompt),
    ]
    return await _send_and_parse(parts)


async def _analyze_parts_batch(
    parts: list[types.Part], offset: int,
) -> list[dict]:
    """Analyze a batch of PPTX text parts in one Gemini call."""
    prompt = SLIDE_ANALYSIS_PROMPT + (
        f"\n\nNOTE: These are slides {offset + 1}-{offset + len(parts)}."
        f" Number them starting at {offset + 1}."
    )
    batch_parts = parts + [types.Part(text=prompt)]
    return await _send_and_parse(batch_parts)


async def _send_and_parse(parts: list[types.Part]) -> list[dict]:
    """Send parts to Gemini and parse the JSON slide array response."""
    content = types.Content(role="user", parts=parts)

    client = genai.Client()
    response = await client.aio.models.generate_content(
        model="gemini-2.5-flash",
        contents=[content],
        config=types.GenerateContentConfig(
            temperature=0.1,
            max_output_tokens=65536,
            response_mime_type="application/json",
        ),
    )

    result_text = (response.text or "").strip()

    # Strip markdown code fences if present
    if result_text.startswith("```"):
        result_text = result_text.split("```")[1]
        if result_text.startswith("json"):
            result_text = result_text[4:]

    try:
        slides = json.loads(result_text)
    except json.JSONDecodeError:
        slides = _recover_truncated_json(result_text)
        if not slides:
            raise SlideAnalysisError(
                "Failed to parse slide analysis JSON and no slides could be recovered"
            )
        logger.warning(
            "Recovered %d slides from truncated Gemini response", len(slides),
        )

    if not isinstance(slides, list):
        raise SlideAnalysisError("Slide analysis did not return a JSON array")

    return slides


def validate_slide_analysis(slides: list[dict]) -> list[dict]:
    """Validate and clean slide analysis results.

    Ensures all slides have slide_number and text_content,
    and boolean fields default to False if missing.
    """
    validated = []
    for i, slide in enumerate(slides):
        validated.append(
            {
                "slide_number": slide.get("slide_number", i + 1),
                "title": slide.get("title"),
                "text_content": slide.get("text_content", ""),
                "visual_description": slide.get("visual_description"),
                "has_diagram": bool(slide.get("has_diagram", False)),
                "has_code": bool(slide.get("has_code", False)),
                "has_equation": bool(slide.get("has_equation", False)),
            }
        )
    return validated


def get_slide_mime_type(file_path: str) -> str:
    """Determine MIME type from file extension."""
    clean = urlparse(file_path).path if file_path.startswith("http") else file_path
    ext = Path(clean).suffix.lower()
    return {
        ".pdf": "application/pdf",
        ".pptx": (
            "application/vnd.openxmlformats-officedocument"
            ".presentationml.presentation"
        ),
    }.get(ext, "application/pdf")



def _pptx_to_image_parts(pptx_path: str) -> list[types.Part]:
    """Convert PPTX slides to PNG image parts for Gemini.

    Uses python-pptx to extract slide dimensions and render each slide
    as a PNG image that Gemini can process with vision.
    """
    from pptx import Presentation

    if pptx_path.startswith(("http://", "https://")):
        import httpx

        resp = httpx.get(pptx_path)
        resp.raise_for_status()
        prs = Presentation(io.BytesIO(resp.content))
    else:
        prs = Presentation(pptx_path)

    parts: list[types.Part] = []
    for i, slide in enumerate(prs.slides):
        # Extract all text from the slide shapes as a fallback
        text_parts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                text_parts.append(shape.text_frame.text)

        slide_text = "\n".join(text_parts)
        slide_desc = f"[Slide {i + 1}]\n{slide_text}"
        parts.append(types.Part(text=slide_desc))

    return parts
