"""Tests for the document text extraction tool."""

from __future__ import annotations

import io
from unittest.mock import MagicMock, patch

import pytest
from docx import Document
from lecturelink_api.tools.document_tools import extract_document_text
from pptx import Presentation
from pptx.util import Inches

# ---------------------------------------------------------------------------
# Helpers to build in-memory test fixtures
# ---------------------------------------------------------------------------

def _make_docx_bytes() -> bytes:
    """Create a minimal DOCX file in memory with headings, a paragraph, and a table."""
    doc = Document()
    doc.add_heading("Syllabus", level=1)
    doc.add_heading("Course Overview", level=2)
    doc.add_paragraph("This is the course description.")

    table = doc.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "Week"
    table.cell(0, 1).text = "Topic"
    table.cell(1, 0).text = "1"
    table.cell(1, 1).text = "Introduction"

    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def _make_pptx_bytes() -> bytes:
    """Create a minimal PPTX file in memory with a title slide and a content slide."""
    prs = Presentation()

    # Slide 1 — title slide
    slide_layout = prs.slide_layouts[0]  # title layout
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Lecture 1"
    slide.placeholders[1].text = "Welcome to the course"

    # Slide 2 — content with a table
    slide_layout2 = prs.slide_layouts[5]  # blank layout
    slide2 = prs.slides.add_slide(slide_layout2)
    txBox = slide2.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    txBox.text_frame.text = "Key concepts"

    # Add a table
    tbl = slide2.shapes.add_table(2, 2, Inches(1), Inches(2.5), Inches(4), Inches(1.5)).table
    tbl.cell(0, 0).text = "Term"
    tbl.cell(0, 1).text = "Definition"
    tbl.cell(1, 0).text = "API"
    tbl.cell(1, 1).text = "Application Programming Interface"

    # Add speaker notes to slide 2
    slide2.notes_slide.notes_text_frame.text = "Remember to explain APIs"

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ---------------------------------------------------------------------------
# Tests
# ---------------------------------------------------------------------------

@pytest.mark.asyncio
async def test_docx_extraction():
    """DOCX extraction should capture headings, paragraphs, and tables."""
    docx_bytes = _make_docx_bytes()

    result = await extract_document_text(
        file_bytes=docx_bytes,
        file_name="syllabus.docx",
        mime_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    )

    assert result["method"] == "docx"
    assert "error" not in result

    text = result["text"]
    assert "# Syllabus" in text
    assert "## Course Overview" in text
    assert "This is the course description." in text
    # Table rendered as markdown
    assert "| Week | Topic |" in text
    assert "| 1 | Introduction |" in text


@pytest.mark.asyncio
async def test_pptx_extraction():
    """PPTX extraction should capture slide titles, body text, tables, and notes."""
    pptx_bytes = _make_pptx_bytes()

    result = await extract_document_text(
        file_bytes=pptx_bytes,
        file_name="lecture.pptx",
        mime_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
    )

    assert result["method"] == "pptx"
    assert result["page_count"] == 2
    assert "error" not in result

    text = result["text"]
    assert "## Slide 1: Lecture 1" in text
    assert "Welcome to the course" in text
    assert "Key concepts" in text
    # Table
    assert "| Term | Definition |" in text
    assert "| API | Application Programming Interface |" in text
    # Speaker notes
    assert "Remember to explain APIs" in text


@pytest.mark.asyncio
async def test_unsupported_mime_type():
    """Unsupported MIME types should return an error dict, not raise."""
    result = await extract_document_text(
        file_bytes=b"not a real file",
        file_name="data.csv",
        mime_type="text/csv",
    )

    assert result["method"] == "unsupported"
    assert result["text"] == ""
    assert "error" in result
    assert "Unsupported file type" in result["error"]


@pytest.mark.asyncio
async def test_pdf_extraction_calls_gemini():
    """PDF extraction should call Gemini and return the response text."""
    mock_response = MagicMock()
    mock_response.text = "Extracted syllabus content from PDF"

    mock_client = MagicMock()
    mock_client.models.generate_content.return_value = mock_response

    mock_settings = MagicMock(GOOGLE_API_KEY="fake-key")
    with (
        patch("lecturelink_api.tools.document_tools.get_settings", return_value=mock_settings),
        patch("lecturelink_api.tools.document_tools.genai.Client", return_value=mock_client),
    ):
        result = await extract_document_text(
            file_bytes=b"%PDF-1.4 fake pdf content",
            file_name="syllabus.pdf",
            mime_type="application/pdf",
        )

    assert result["method"] == "gemini-pdf"
    assert result["text"] == "Extracted syllabus content from PDF"
    assert "error" not in result

    # Verify the Gemini client was called with the right model
    call_kwargs = mock_client.models.generate_content.call_args
    assert call_kwargs.kwargs["model"] == "gemini-2.5-flash"


@pytest.mark.asyncio
async def test_pdf_extraction_handles_gemini_error():
    """If Gemini fails, we should get an error dict instead of an exception."""
    mock_client = MagicMock()
    mock_client.models.generate_content.side_effect = RuntimeError("API quota exceeded")

    mock_settings = MagicMock(GOOGLE_API_KEY="fake-key")
    with (
        patch("lecturelink_api.tools.document_tools.get_settings", return_value=mock_settings),
        patch("lecturelink_api.tools.document_tools.genai.Client", return_value=mock_client),
    ):
        result = await extract_document_text(
            file_bytes=b"%PDF-1.4 fake pdf content",
            file_name="syllabus.pdf",
            mime_type="application/pdf",
        )

    assert result["method"] == "gemini-pdf"
    assert result["text"] == ""
    assert "error" in result
    assert "API quota exceeded" in result["error"]


# ---------------------------------------------------------------------------
# Text box extraction tests
# ---------------------------------------------------------------------------


class TestTextboxExtraction:
    def test_finds_text_in_txbxContent(self):
        from lecturelink_api.tools.document_tools import _extract_textbox_content
        from lxml import etree

        xml_str = (
            '<mc:AlternateContent'
            '  xmlns:mc="http://schemas.openxmlformats.org/markup-compatibility/2006"'
            '  xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main"'
            '  xmlns:wps="http://schemas.microsoft.com/office/word/2010/wordprocessingShape">'
            '  <mc:Choice Requires="wps">'
            '    <w:drawing><wps:wsp><wps:txbx><w:txbxContent>'
            '      <w:p><w:r><w:t>MMAI 5090 3.00</w:t></w:r></w:p>'
            '      <w:p><w:r><w:t>Business Applications of AI II</w:t></w:r></w:p>'
            '    </w:txbxContent></wps:txbx></wps:wsp></w:drawing>'
            '  </mc:Choice>'
            '</mc:AlternateContent>'
        )
        element = etree.fromstring(xml_str)
        texts = _extract_textbox_content(element)
        assert len(texts) == 2
        assert texts[0] == "MMAI 5090 3.00"
        assert texts[1] == "Business Applications of AI II"

    def test_skips_empty_paragraphs(self):
        from lecturelink_api.tools.document_tools import _extract_textbox_content
        from lxml import etree

        xml_str = (
            '<mc:AlternateContent'
            '  xmlns:mc="http://schemas.openxmlformats.org/markup-compatibility/2006"'
            '  xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main"'
            '  xmlns:wps="http://schemas.microsoft.com/office/word/2010/wordprocessingShape">'
            '  <mc:Choice Requires="wps">'
            '    <wps:wsp><wps:txbx><w:txbxContent>'
            '      <w:p></w:p>'
            '      <w:p><w:r><w:t>Content</w:t></w:r></w:p>'
            '    </w:txbxContent></wps:txbx></wps:wsp>'
            '  </mc:Choice>'
            '</mc:AlternateContent>'
        )
        element = etree.fromstring(xml_str)
        texts = _extract_textbox_content(element)
        assert texts == ["Content"]

    def test_no_textboxes_returns_empty(self):
        from lecturelink_api.tools.document_tools import _extract_textbox_content
        from lxml import etree

        xml_str = '<w:p xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main"><w:r><w:t>Text</w:t></w:r></w:p>'
        element = etree.fromstring(xml_str)
        assert _extract_textbox_content(element) == []

    def test_inline_textbox_inside_paragraph(self):
        """Text boxes embedded inside w:p elements (inline drawings) should be extracted."""
        from lecturelink_api.tools.document_tools import _extract_textbox_content
        from lxml import etree

        # Simulate a paragraph containing an inline drawing with a text box
        xml_str = (
            '<w:p xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main"'
            '     xmlns:wps="http://schemas.microsoft.com/office/word/2010/wordprocessingShape">'
            '  <w:r><w:t>Regular text</w:t></w:r>'
            '  <w:r><w:drawing><wps:wsp><wps:txbx><w:txbxContent>'
            '    <w:p><w:r><w:t>MMAI 5090 3.00</w:t></w:r></w:p>'
            '  </w:txbxContent></wps:txbx></wps:wsp></w:drawing></w:r>'
            '</w:p>'
        )
        element = etree.fromstring(xml_str)
        texts = _extract_textbox_content(element)
        assert texts == ["MMAI 5090 3.00"]

    def test_multi_run_concatenation(self):
        from lecturelink_api.tools.document_tools import _extract_textbox_content
        from lxml import etree

        xml_str = (
            '<mc:AlternateContent'
            '  xmlns:mc="http://schemas.openxmlformats.org/markup-compatibility/2006"'
            '  xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main"'
            '  xmlns:wps="http://schemas.microsoft.com/office/word/2010/wordprocessingShape">'
            '  <mc:Choice Requires="wps">'
            '    <wps:wsp><wps:txbx><w:txbxContent>'
            '      <w:p><w:r><w:t>MMAI </w:t></w:r><w:r><w:t>5090</w:t></w:r></w:p>'
            '    </w:txbxContent></wps:txbx></wps:wsp>'
            '  </mc:Choice>'
            '</mc:AlternateContent>'
        )
        element = etree.fromstring(xml_str)
        texts = _extract_textbox_content(element)
        assert texts == ["MMAI 5090"]


# ---------------------------------------------------------------------------
# Heading detection tests
# ---------------------------------------------------------------------------


class TestHeadingDetection:
    def test_standard_headings(self):
        from lecturelink_api.tools.document_tools import _get_heading_prefix

        assert _get_heading_prefix("Heading 1") == "# "
        assert _get_heading_prefix("Heading 2") == "## "
        assert _get_heading_prefix("Heading 3") == "### "
        assert _get_heading_prefix("Heading 4") == "#### "

    def test_custom_heading_styles(self):
        from lecturelink_api.tools.document_tools import _get_heading_prefix

        assert _get_heading_prefix("Page 1 Heading 1") == "# "
        assert _get_heading_prefix("Custom Heading 2") == "## "

    def test_title_subtitle(self):
        from lecturelink_api.tools.document_tools import _get_heading_prefix

        assert _get_heading_prefix("Title") == "# "
        assert _get_heading_prefix("Subtitle") == "## "
        assert _get_heading_prefix("Course Title") == "# "

    def test_toc_skipped(self):
        from lecturelink_api.tools.document_tools import _get_heading_prefix

        assert _get_heading_prefix("toc 1") is None
        assert _get_heading_prefix("TOC Heading") is None

    def test_normal_returns_empty(self):
        from lecturelink_api.tools.document_tools import _get_heading_prefix

        assert _get_heading_prefix("Normal") == ""
        assert _get_heading_prefix("") == ""

    def test_heading_without_number_defaults_to_h2(self):
        from lecturelink_api.tools.document_tools import _get_heading_prefix

        assert _get_heading_prefix("Custom Heading") == "## "


# ---------------------------------------------------------------------------
# Merged cell dedup tests
# ---------------------------------------------------------------------------


class TestMergedCellDedup:
    @pytest.mark.asyncio
    async def test_merged_cells_deduplicated(self):
        doc = Document()
        table = doc.add_table(rows=2, cols=3)
        table.cell(0, 0).text = "DATE"
        table.cell(0, 1).text = "DATE"
        table.cell(0, 2).text = "TOPIC"
        table.cell(1, 0).text = "Week 1"
        table.cell(1, 1).text = "Week 1"
        table.cell(1, 2).text = "Introduction"

        buf = io.BytesIO()
        doc.save(buf)

        result = await extract_document_text(
            file_bytes=buf.getvalue(),
            file_name="test.docx",
            mime_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        )
        text = result["text"]
        assert "| DATE | TOPIC |" in text
        assert "| DATE | DATE |" not in text
        assert "| Week 1 | Introduction |" in text


# ---------------------------------------------------------------------------
# List paragraph tests
# ---------------------------------------------------------------------------


class TestListParagraphs:
    @pytest.mark.asyncio
    async def test_list_paragraphs_get_bullet_prefix(self):
        doc = Document()
        doc.add_heading("Learning Outcomes", level=1)
        for text in ["Understand AI", "Apply ML models"]:
            para = doc.add_paragraph(text)
            para.style = doc.styles["List Paragraph"]

        buf = io.BytesIO()
        doc.save(buf)

        result = await extract_document_text(
            file_bytes=buf.getvalue(),
            file_name="test.docx",
            mime_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        )
        text = result["text"]
        assert "- Understand AI" in text
        assert "- Apply ML models" in text


# ---------------------------------------------------------------------------
# Regression tests
# ---------------------------------------------------------------------------


class TestDocxExtractionRegression:
    @pytest.mark.asyncio
    async def test_basic_extraction_still_works(self):
        docx_bytes = _make_docx_bytes()
        result = await extract_document_text(
            file_bytes=docx_bytes,
            file_name="syllabus.docx",
            mime_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        )
        assert result["method"] == "docx"
        assert "error" not in result
        assert "# Syllabus" in result["text"]
        assert "## Course Overview" in result["text"]
        assert "This is the course description." in result["text"]
        assert "| Week | Topic |" in result["text"]
